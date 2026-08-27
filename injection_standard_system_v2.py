# -*- coding: utf-8 -*-
"""
사출 표준서 자동생성 시스템 v2
====================================================
[문서 자동화 워크플로우]
  STEP 1. 문서 카테고리 선정 (사출 TO 표준서 / 수지관리 표준서 / 신규 카테고리)
  STEP 2. 카테고리별 항목 선정 (기본항목 + AI제안 + 사용자 직접입력/노하우)
  STEP 3. 과거 표준서 참고 → 형식/양식(폰트, 표스타일, 색상) 자동 매칭
  STEP 4. v1 문서 검토 후 "집중관리 포인트" 지정 → 최종본 확정
  STEP 5. 워드(.docx) / PPT(.pptx) / 엑셀 체크리스트(.xlsx) 일괄 다운로드

운영 방식: 하이브리드
  - AI(Gemini): 항목 제안/삭제 제안, 서술 보완만 담당 (판단은 사용자)
  - 로컬 로직: 표/수치/스타일 조립, 문서 생성은 100% 결정론적 코드로 처리

배포: Streamlit Community Cloud (기존 IMS/AUTO-DESIGN 앱과 동일 패턴)
secrets.toml:
    GEMINI_API_KEY = "..."
"""

import json
import io
from datetime import date

import streamlit as st
import google.generativeai as genai
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PPTRGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.worksheet.datavalidation import DataValidation

st.set_page_config(page_title="사출 표준서 자동생성 시스템", layout="wide")
GEMINI_MODEL_NAME = "gemini-2.5-flash"


def get_model():
    api_key = st.secrets.get("GEMINI_API_KEY", "")
    if not api_key:
        st.error("st.secrets['GEMINI_API_KEY'] 가 설정되어 있지 않습니다.")
        st.stop()
    genai.configure(api_key=api_key)
    return genai.GenerativeModel(GEMINI_MODEL_NAME)


def _parse_json_array(text: str) -> list:
    cleaned = text.strip().replace("```json", "").replace("```", "").strip()
    try:
        data = json.loads(cleaned)
        return data if isinstance(data, list) else []
    except json.JSONDecodeError:
        s, e = cleaned.find("["), cleaned.rfind("]")
        if s != -1 and e != -1:
            try:
                return json.loads(cleaned[s:e + 1])
            except json.JSONDecodeError:
                pass
    return []


# ============================================================
# STEP 1. 카테고리 템플릿 (기본항목은 AI가 아니라 고정 로직 — 안정성 확보)
# ============================================================
CATEGORY_TEMPLATES = {
    "사출 TO 표준서": {
        "desc": "신규/변경 금형의 Try-Out 단계에서 성형조건을 확정하고 품질기준을 문서화하는 표준서",
        "sections": ["성형 조건", "온도 관리", "치수/중량", "설비", "품질"],
        "default_items": [
            {"category": "성형 조건", "item": "사출압력(1차)", "criteria": "설정치 대비 ±5% 이내", "basis": "충전 불균일 및 버(Burr) 방지", "checklist": True},
            {"category": "성형 조건", "item": "보압 / 보압시간", "criteria": "게이트 씰(Seal) 시점 이후 0.5~1.0초 여유", "basis": "싱크마크·수축 방지", "checklist": True},
            {"category": "성형 조건", "item": "사출속도", "criteria": "웰드라인/플로우마크 미발생 구간", "basis": "외관 품질 확보", "checklist": True},
            {"category": "성형 조건", "item": "냉각시간", "criteria": "이젝션 시 변형 없는 최소시간", "basis": "치수 안정성, 생산성 균형", "checklist": True},
            {"category": "온도 관리", "item": "실린더 온도(Zone별)", "criteria": "재료 TDS 권장범위 이내", "basis": "재료 열화/분해 방지", "checklist": True},
            {"category": "온도 관리", "item": "금형 온도", "criteria": "재료 권장 금형온도 ±5℃", "basis": "결정화/전사성 확보", "checklist": True},
            {"category": "치수/중량", "item": "샷중량(Shot Weight)", "criteria": "기준중량 대비 ±1% 이내", "basis": "충전 재현성 확인", "checklist": True},
            {"category": "치수/중량", "item": "쿠션량(Cushion)", "criteria": "설정 쿠션 대비 변동 0.5mm 이내", "basis": "역류 방지, 보압 전달 확인", "checklist": True},
            {"category": "설비", "item": "형체력", "criteria": "투영면적 기준 산출값 이내", "basis": "버/플래시 방지", "checklist": False},
            {"category": "품질", "item": "외관 검사", "criteria": "웰드라인/플로우마크/실버 등 육안 기준 이내", "basis": "고객 외관 스펙 준수", "checklist": True},
            {"category": "품질", "item": "치수 검사(주요부)", "criteria": "도면 공차 이내(Cpk 관리항목은 별도 표기)", "basis": "조립성/기능성 확보", "checklist": True},
        ],
    },
    "플라스틱 수지(재료) 관리 표준서": {
        "desc": "원재료 입고부터 보관·건조·이력관리까지 사출 재료 품질을 관리하는 표준서",
        "sections": ["입고 관리", "보관 조건", "건조 조건", "이력 관리", "품질 검사"],
        "default_items": [
            {"category": "입고 관리", "item": "입고 검사 성적서 확인", "criteria": "COA(성적서) 로트번호와 실물 라벨 일치", "basis": "오사용/오입고 방지", "checklist": True},
            {"category": "입고 관리", "item": "포장 상태 확인", "criteria": "방습포장 파손·개봉 흔적 없음", "basis": "흡습에 의한 가수분해 방지", "checklist": True},
            {"category": "보관 조건", "item": "보관 온습도", "criteria": "온도 25℃ 이하, 습도 60%RH 이하", "basis": "흡습 및 재료 열화 방지", "checklist": True},
            {"category": "보관 조건", "item": "선입선출(FIFO)", "criteria": "입고일 기준 오래된 로트 우선 사용", "basis": "장기 보관에 따른 물성 저하 방지", "checklist": True},
            {"category": "건조 조건", "item": "건조 온도/시간", "criteria": "재료 TDS 권장 건조조건 준수", "basis": "수분에 의한 실버/기포 불량 방지", "checklist": True},
            {"category": "건조 조건", "item": "건조 후 수분율", "criteria": "재료별 기준 수분율 이내(예: PA66 0.2% 이하)", "basis": "가수분해로 인한 물성 저하 방지", "checklist": True},
            {"category": "이력 관리", "item": "로트별 사용 이력 기록", "criteria": "로트번호-생산Lot 매칭 기록 유지", "basis": "품질 이슈 발생 시 추적성 확보", "checklist": True},
            {"category": "품질 검사", "item": "MFI(용융지수) 확인", "criteria": "재료 스펙 대비 ±10% 이내", "basis": "성형성 및 물성 재현성 확인", "checklist": False},
        ],
    },
}
CUSTOM_LABEL = "신규 카테고리 직접 입력"


# ============================================================
# AI 함수 - 신규 카테고리 초기 항목 생성 / 추가제안 / 삭제제안
# ============================================================
def ai_generate_category_items(category_name: str, purpose: str) -> list:
    model = get_model()
    prompt = f"""
너는 30년 경력의 사출성형 제조 현장 전문가야. "{category_name}"이라는 표준서를 새로 만들려고 해.
목적/설명: {purpose}

이 표준서에 들어가야 할 섹션(카테고리) 4~6개와, 각 섹션당 항목을 2~4개씩 제안해줘.
각 항목에는 판단기준(구체적 수치/조건)과 판단근거(왜 그 기준인지)를 포함해줘.

반드시 아래 JSON 배열 형식으로만 응답하고 그 외 텍스트는 포함하지 마.
[
  {{"category": "섹션명", "item": "항목명", "criteria": "판단기준", "basis": "판단근거", "checklist": true}}
]
"""
    resp = model.generate_content(prompt)
    return _parse_json_array(resp.text)


def ai_suggest_additions(context: dict, current_items: list) -> list:
    model = get_model()
    names = [it["item"] for it in current_items]
    prompt = f"""
너는 30년 경력의 사출성형 공정 전문가야. 문서 카테고리: "{context.get('category_name')}"
[배경정보] 제품명: {context.get('product','미기재')} / 재질: {context.get('material','미기재')} /
금형정보: {context.get('mold_info','미기재')} / 특이사항: {context.get('note','미기재')}

[현재 항목 목록] {json.dumps(names, ensure_ascii=False)}

배경정보를 고려해 빠진 "추가 항목"을 최대 6개 제안해줘. 중복 금지.
JSON 배열만 응답:
[
  {{"category": "섹션명(기존 섹션 중 하나 또는 신규)", "item": "항목명", "criteria": "판단기준",
    "basis": "판단근거", "checklist": true, "reason_to_add": "제안 사유"}}
]
"""
    resp = model.generate_content(prompt)
    return _parse_json_array(resp.text)


def ai_suggest_removals(context: dict, current_items: list) -> list:
    model = get_model()
    prompt = f"""
너는 30년 경력의 사출성형 공정 전문가야. 문서 카테고리: "{context.get('category_name')}"
[배경정보] 제품명: {context.get('product','미기재')} / 재질: {context.get('material','미기재')} /
금형정보: {context.get('mold_info','미기재')} / 특이사항: {context.get('note','미기재')}

[현재 항목] {json.dumps([{"item": it["item"], "category": it["category"]} for it in current_items], ensure_ascii=False)}

실효성이 낮거나 중복되는 "삭제 검토 항목"을 최대 4개만 골라줘(없으면 빈 배열).
JSON 배열만 응답:
[{{"item": "정확한 항목명", "reason_to_remove": "삭제 검토 사유"}}]
"""
    resp = model.generate_content(prompt)
    return _parse_json_array(resp.text)


def ai_suggest_basis(context: dict, items: list) -> list:
    """현재 항목 전체를 검토해서 판단근거가 부실하거나 비어있는 항목을 보완 제안.
    이미 근거가 있는 항목도 더 구체적이고 설득력 있게 다듬어 제안한다."""
    model = get_model()
    payload = [{"item": it["item"], "category": it["category"], "criteria": it["criteria"],
                "current_basis": it.get("basis", "")} for it in items]
    prompt = f"""
너는 30년 경력의 사출성형 품질/공정 전문가야. 문서 카테고리: "{context.get('category_name')}"
[배경정보] 제품명: {context.get('product','미기재')} / 재질: {context.get('material','미기재')} /
금형정보: {context.get('mold_info','미기재')} / 특이사항: {context.get('note','미기재')}

아래 항목들의 "판단근거"를 검토해줘. 근거가 비어있거나 막연한 항목은 구체적인 공정/품질 원리에
근거해 새로 작성하고, 이미 근거가 충실한 항목은 그대로 두거나 더 다듬어도 좋아.
근거는 1~2문장, 왜 그 판단기준이 타당한지 공정/재료/품질 관점에서 설명해줘.

[항목 목록]
{json.dumps(payload, ensure_ascii=False, indent=2)}

반드시 아래 JSON 배열 형식으로만 응답하고 그 외 텍스트는 포함하지 마. 근거를 새로 제안하는
항목만 포함해도 되고, current_basis와 동일하면 포함하지 마.
[{{"item": "정확한 항목명", "suggested_basis": "보완된 판단근거"}}]
"""
    resp = model.generate_content(prompt)
    return _parse_json_array(resp.text)


def ai_final_review(context: dict, items: list) -> str:
    model = get_model()
    prompt = f"""
너는 30년 경력의 사출성형 품질 전문가야. 아래는 "{context.get('category_name')}" 표준서의
최종 항목 목록이야. 전체적으로 빠진 게 없는지, 항목 간 중복/모순은 없는지 검토하고
3~5줄로 간단한 검토 코멘트를 줘 (항목을 나열하지 말고 총평 위주로).

[항목 목록]
{json.dumps([{"category": it["category"], "item": it["item"]} for it in items], ensure_ascii=False, indent=2)}
"""
    resp = model.generate_content(prompt)
    return resp.text.strip()


# ============================================================
# STEP 3. 참고 서식(과거 표준서) 스타일 분석
# ============================================================
def analyze_reference_docx(file_bytes: bytes) -> dict:
    """업로드된 과거 표준서에서 헤딩폰트/색상/표스타일을 추출해 스타일 프로파일 생성.
    실패하거나 정보가 부족하면 None 값으로 채워 기본 스타일로 폴백."""
    profile = {"heading_font": None, "heading_size": None, "heading_color": None,
               "body_font": None, "body_size": None, "table_style": None}
    try:
        doc = Document(io.BytesIO(file_bytes))
        for p in doc.paragraphs:
            if p.runs:
                r = p.runs[0]
                if r.bold and r.font.size and r.font.size.pt >= 13:
                    profile["heading_font"] = r.font.name
                    profile["heading_size"] = r.font.size.pt
                    if r.font.color and r.font.color.rgb:
                        profile["heading_color"] = str(r.font.color.rgb)
                    break
        for p in doc.paragraphs:
            if p.runs and not p.runs[0].bold and p.runs[0].text.strip():
                r = p.runs[0]
                profile["body_font"] = r.font.name
                if r.font.size:
                    profile["body_size"] = r.font.size.pt
                break
        if doc.tables:
            try:
                profile["table_style"] = doc.tables[0].style.name
            except Exception:
                profile["table_style"] = None
    except Exception as e:
        st.warning(f"참고 파일 분석 중 일부 정보를 읽지 못했습니다: {e}")
    return profile


DEFAULT_STYLE = {
    "heading_font": "맑은 고딕", "heading_size": 13, "heading_color": "1F4E79",
    "body_font": "맑은 고딕", "body_size": 10, "table_style": "Light Grid Accent 1",
}


def merged_style(profile: dict) -> dict:
    """참고파일에서 추출된 값이 있으면 우선 사용, 없으면 기본값으로 채움."""
    return {k: (profile.get(k) if profile and profile.get(k) else v) for k, v in DEFAULT_STYLE.items()}


# ============================================================
# 세션 상태 초기화
# ============================================================
def init_state():
    ss = st.session_state
    ss.setdefault("stage", 1)
    ss.setdefault("category_name", "")
    ss.setdefault("category_purpose", "")
    ss.setdefault("items", [])
    ss.setdefault("next_id", 0)
    ss.setdefault("ai_add_suggestions", [])
    ss.setdefault("ai_remove_suggestions", [])
    ss.setdefault("ai_basis_suggestions", [])
    ss.setdefault("style_profile", {})
    ss.setdefault("context", {})
    ss.setdefault("final_review_text", "")


init_state()


def new_id():
    st.session_state.next_id += 1
    return st.session_state.next_id - 1


def goto(stage):
    st.session_state.stage = stage
    st.rerun()


# ============================================================
# 문서 생성 함수들
# ============================================================
def build_docx(context: dict, items: list, style: dict) -> bytes:
    doc = Document()
    heading_rgb = RGBColor.from_string(style["heading_color"]) if style["heading_color"] else RGBColor(0x1F, 0x4E, 0x79)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run(context.get("category_name", "표준서"))
    r.bold = True
    r.font.size = Pt(20)
    r.font.name = style["heading_font"]

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run(f"제품명: {context.get('product','')}").font.size = Pt(12)

    info_table = doc.add_table(rows=2, cols=4)
    try:
        info_table.style = style["table_style"]
    except Exception:
        info_table.style = "Light Grid Accent 1"
    hdr = ["재질", "금형 정보", "작성일", "개정번호"]
    vals = [context.get("material", ""), context.get("mold_info", ""), str(date.today()), context.get("rev", "Rev.0")]
    for i, (h, v) in enumerate(zip(hdr, vals)):
        info_table.cell(0, i).text = h
        info_table.cell(1, i).text = v
        for p in info_table.cell(0, i).paragraphs:
            for run in p.runs:
                run.bold = True

    doc.add_paragraph("")
    if context.get("note"):
        p = doc.add_paragraph()
        p.add_run("특이사항: ").bold = True
        p.add_run(context["note"])

    # 집중관리 포인트 요약 섹션
    focus_items = [it for it in items if it.get("focus")]
    if focus_items:
        doc.add_paragraph("")
        fh = doc.add_paragraph()
        fr = fh.add_run("집중관리 포인트")
        fr.bold = True
        fr.font.size = Pt(13)
        fr.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
        for it in focus_items:
            bp = doc.add_paragraph(style=None)
            bp.add_run(f"• {it['item']} — {it['criteria']}").bold = True

    doc.add_paragraph("")

    categories = []
    for it in items:
        if it["category"] not in categories:
            categories.append(it["category"])

    for cat in categories:
        heading = doc.add_paragraph()
        hr = heading.add_run(f"■ {cat}")
        hr.bold = True
        hr.font.size = Pt(13)
        hr.font.color.rgb = heading_rgb
        hr.font.name = style["heading_font"]

        table = doc.add_table(rows=1, cols=6)
        try:
            table.style = style["table_style"]
        except Exception:
            table.style = "Light Grid Accent 1"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        widths = [Cm(3.0), Cm(4.0), Cm(4.5), Cm(1.5), Cm(3.5), Cm(1.5)]
        headers = ["항목", "판단기준", "판단 근거", "체크", "현장 노하우/비고", "집중"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.width = widths[i]
            for p in cell.paragraphs:
                for run in p.runs:
                    run.bold = True
                    run.font.size = Pt(style["body_size"])

        for it in [x for x in items if x["category"] == cat]:
            row = table.add_row()
            values = [it["item"], it["criteria"], it["basis"],
                      "[  ]" if it.get("checklist") else "-",
                      it.get("user_note", "") or it.get("note", ""),
                      "FOCUS" if it.get("focus") else ""]
            for i, v in enumerate(values):
                row.cells[i].text = v
                row.cells[i].width = widths[i]
                for p in row.cells[i].paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(style["body_size"])
                        run.font.name = style["body_font"]
            if it.get("focus"):
                for cell in row.cells:
                    shd = cell._tc.get_or_add_tcPr()
                    from docx.oxml.ns import qn as _qn
                    from docx.oxml import OxmlElement
                    shade = OxmlElement("w:shd")
                    shade.set(_qn("w:fill"), "FCE4E4")
                    shd.append(shade)

        doc.add_paragraph("")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def build_pptx(context: dict, items: list, style: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    heading_rgb = PPTRGBColor.from_string(style["heading_color"]) if style["heading_color"] else PPTRGBColor(0x1F, 0x4E, 0x79)

    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2))
    tf = box.text_frame
    tf.text = context.get("category_name", "표준서")
    tf.paragraphs[0].font.size = PPt(36)
    tf.paragraphs[0].font.bold = True

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.0))
    sub_box.text_frame.text = f"제품명: {context.get('product','')}  |  재질: {context.get('material','')}  |  {date.today()}"
    sub_box.text_frame.paragraphs[0].font.size = PPt(16)

    # 집중관리 포인트 요약 슬라이드
    focus_items = [it for it in items if it.get("focus")]
    if focus_items:
        fslide = prs.slides.add_slide(blank)
        ftitle = fslide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        ftitle.text_frame.text = "집중관리 포인트"
        ftitle.text_frame.paragraphs[0].font.size = PPt(28)
        ftitle.text_frame.paragraphs[0].font.bold = True
        ftitle.text_frame.paragraphs[0].font.color.rgb = PPTRGBColor(0xC0, 0x00, 0x00)
        body = fslide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(5.5))
        btf = body.text_frame
        for i, it in enumerate(focus_items):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.text = f"• {it['item']}  —  {it['criteria']}"
            p.font.size = PPt(18)

    categories = []
    for it in items:
        if it["category"] not in categories:
            categories.append(it["category"])

    for cat in categories:
        cat_items = [x for x in items if x["category"] == cat]
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_box.text_frame.text = f"■ {cat}"
        title_box.text_frame.paragraphs[0].font.size = PPt(26)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = heading_rgb

        rows, cols = len(cat_items) + 1, 4
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5 * rows))
        table = table_shape.table
        for i, w in enumerate([Inches(2.6), Inches(3.4), Inches(3.9), Inches(2.4)]):
            table.columns[i].width = w
        headers = ["항목", "판단기준", "판단 근거", "현장 노하우/비고"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = PPt(13)
            cell.fill.solid()
            cell.fill.fore_color.rgb = heading_rgb
            cell.text_frame.paragraphs[0].font.color.rgb = PPTRGBColor(0xFF, 0xFF, 0xFF)
        for r, it in enumerate(cat_items, start=1):
            vals = [it["item"], it["criteria"], it["basis"], it.get("user_note", "") or it.get("note", "")]
            for c, v in enumerate(vals):
                cell = table.cell(r, c)
                cell.text = v
                cell.text_frame.paragraphs[0].font.size = PPt(11)
                if it.get("focus"):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PPTRGBColor(0xFC, 0xE4, 0xE4)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def build_xlsx(context: dict, items: list) -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "체크리스트"

    header_fill = PatternFill("solid", fgColor="1F4E79")
    focus_fill = PatternFill("solid", fgColor="FCE4E4")
    thin = Side(style="thin", color="D9D9D9")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    ws["A1"] = context.get("category_name", "표준서")
    ws["A1"].font = Font(name="Arial", size=16, bold=True)
    ws.merge_cells("A1:G1")
    ws["A2"] = f"제품명: {context.get('product','')}   재질: {context.get('material','')}   작성일: {date.today()}"
    ws["A2"].font = Font(name="Arial", size=10, italic=True)
    ws.merge_cells("A2:G2")

    headers = ["카테고리", "항목", "판단기준", "판단 근거", "체크(Pass/Fail/N-A)", "현장 노하우/비고", "집중포인트"]
    header_row = 4
    for i, h in enumerate(headers, start=1):
        c = ws.cell(row=header_row, column=i, value=h)
        c.font = Font(name="Arial", size=11, bold=True, color="FFFFFF")
        c.fill = header_fill
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border = border
    ws.freeze_panes = f"A{header_row + 1}"

    dv = DataValidation(type="list", formula1='"Pass,Fail,N/A"', allow_blank=True)
    ws.add_data_validation(dv)

    r = header_row + 1
    for it in items:
        row_vals = [it["category"], it["item"], it["criteria"], it["basis"], "",
                    it.get("user_note", "") or it.get("note", ""), "FOCUS" if it.get("focus") else ""]
        for c_idx, v in enumerate(row_vals, start=1):
            cell = ws.cell(row=r, column=c_idx, value=v)
            cell.font = Font(name="Arial", size=10)
            cell.alignment = Alignment(vertical="center", wrap_text=True)
            cell.border = border
            if it.get("focus"):
                cell.fill = focus_fill
        dv.add(ws.cell(row=r, column=5))
        r += 1

    widths = [16, 22, 30, 32, 18, 26, 10]
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[openpyxl.utils.get_column_letter(i)].width = w

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================
# UI - 상단 진행 표시
# ============================================================
STAGE_NAMES = ["카테고리 선정", "항목 선정", "서식 매칭", "집중포인트/최종본", "다운로드"]
st.title("사출 표준서 자동생성 시스템")
prog_cols = st.columns(5)
for i, name in enumerate(STAGE_NAMES, start=1):
    with prog_cols[i - 1]:
        marker = "[진행중]" if i == st.session_state.stage else ("[완료]" if i < st.session_state.stage else "[대기]")
        st.markdown(f"{marker} **{i}. {name}**")
st.divider()

# ============================================================
# STAGE 1. 카테고리 선정
# ============================================================
if st.session_state.stage == 1:
    st.subheader("STEP 1. 문서 카테고리 선정")
    options = list(CATEGORY_TEMPLATES.keys()) + [CUSTOM_LABEL]
    choice = st.selectbox("표준서 카테고리를 선택하세요", options)

    if choice == CUSTOM_LABEL:
        custom_name = st.text_input("신규 카테고리 이름", placeholder="예: 사출 금형 예방보전 표준서")
        purpose = st.text_area("이 표준서의 목적/범위를 간단히 설명해주세요",
                                placeholder="예: 금형 정기 예방보전 주기와 점검항목을 관리하기 위한 표준서")
        if st.button("다음 단계 →", type="primary"):
            if not custom_name.strip():
                st.warning("카테고리 이름을 입력해주세요.")
            else:
                with st.spinner("AI가 카테고리 구조를 설계 중..."):
                    generated = ai_generate_category_items(custom_name, purpose)
                st.session_state.category_name = custom_name
                st.session_state.category_purpose = purpose
                st.session_state.items = [dict(it, id=new_id(), focus=False) for it in generated]
                goto(2)
    else:
        tpl = CATEGORY_TEMPLATES[choice]
        st.info(tpl["desc"])
        st.caption("포함 섹션: " + " / ".join(tpl["sections"]))
        if st.button("다음 단계 →", type="primary"):
            st.session_state.category_name = choice
            st.session_state.category_purpose = tpl["desc"]
            st.session_state.items = [dict(it, id=new_id(), focus=False) for it in tpl["default_items"]]
            goto(2)

# ============================================================
# STAGE 2. 항목 선정 (AI + 사용자)
# ============================================================
elif st.session_state.stage == 2:
    st.subheader(f"STEP 2. 항목 선정 — {st.session_state.category_name}")

    with st.expander("제품/공정 배경정보 입력 (AI 제안 정확도에 영향)", expanded=True):
        c1, c2 = st.columns(2)
        product = c1.text_input("제품명", value=st.session_state.context.get("product", ""))
        material = c2.text_input("재질", value=st.session_state.context.get("material", ""))
        mold_info = c1.text_input("금형 정보", value=st.session_state.context.get("mold_info", ""))
        note = c2.text_input("특이사항", value=st.session_state.context.get("note", ""))
        rev = c1.text_input("개정번호", value=st.session_state.context.get("rev", "Rev.0"))
        st.session_state.context = {"category_name": st.session_state.category_name,
                                     "product": product, "material": material,
                                     "mold_info": mold_info, "note": note, "rev": rev}

    colA, colB, colC = st.columns(3)
    if colA.button("AI 추가 항목 제안받기", use_container_width=True):
        with st.spinner("분석 중..."):
            st.session_state.ai_add_suggestions = ai_suggest_additions(st.session_state.context, st.session_state.items)
    if colB.button("AI 삭제 검토 항목 제안받기", use_container_width=True):
        with st.spinner("분석 중..."):
            st.session_state.ai_remove_suggestions = ai_suggest_removals(st.session_state.context, st.session_state.items)
    if colC.button("AI 판단근거 보완 제안받기", use_container_width=True):
        with st.spinner("분석 중..."):
            st.session_state.ai_basis_suggestions = ai_suggest_basis(st.session_state.context, st.session_state.items)

    if st.session_state.ai_add_suggestions:
        st.markdown("**AI 추가 항목 제안** (체크한 것만 반영)")
        to_add = []
        for i, sug in enumerate(st.session_state.ai_add_suggestions):
            c1, c2 = st.columns([0.06, 0.94])
            checked = c1.checkbox("", key=f"add_{i}")
            c2.markdown(f"**[{sug.get('category','기타')}] {sug.get('item','')}** — {sug.get('criteria','')}\n\n"
                        f"근거: {sug.get('basis','')} · _{sug.get('reason_to_add','')}_")
            if checked:
                to_add.append(sug)
        if st.button("선택 항목 반영"):
            for sug in to_add:
                st.session_state.items.append({"id": new_id(), "category": sug.get("category", "기타"),
                                                 "item": sug.get("item", ""), "criteria": sug.get("criteria", ""),
                                                 "basis": sug.get("basis", ""), "checklist": sug.get("checklist", True),
                                                 "user_note": "", "focus": False})
            st.session_state.ai_add_suggestions = []
            st.rerun()

    if st.session_state.ai_remove_suggestions:
        st.markdown("**AI 삭제 검토 항목** (체크한 것만 삭제)")
        to_remove = []
        for i, sug in enumerate(st.session_state.ai_remove_suggestions):
            c1, c2 = st.columns([0.06, 0.94])
            checked = c1.checkbox("", key=f"rm_{i}")
            c2.markdown(f"**{sug.get('item','')}** — _{sug.get('reason_to_remove','')}_")
            if checked:
                to_remove.append(sug.get("item"))
        if st.button("선택 항목 삭제"):
            st.session_state.items = [it for it in st.session_state.items if it["item"] not in to_remove]
            st.session_state.ai_remove_suggestions = []
            st.rerun()

    if st.session_state.ai_basis_suggestions:
        st.markdown("**AI 판단근거 보완 제안** (체크한 것만 반영)")
        name_to_item = {it["item"]: it for it in st.session_state.items}
        to_apply = []
        for i, sug in enumerate(st.session_state.ai_basis_suggestions):
            target = name_to_item.get(sug.get("item"))
            if not target:
                continue
            c1, c2 = st.columns([0.06, 0.94])
            checked = c1.checkbox("", key=f"basis_{i}")
            c2.markdown(f"**{sug.get('item','')}**\n\n"
                        f"기존 근거: _{target.get('basis','') or '(없음)'}_\n\n"
                        f"제안 근거: {sug.get('suggested_basis','')}")
            if checked:
                to_apply.append(sug)
        if st.button("선택 항목 근거 반영"):
            for sug in to_apply:
                target = name_to_item.get(sug.get("item"))
                if target:
                    target["basis"] = sug.get("suggested_basis", target.get("basis", ""))
            st.session_state.ai_basis_suggestions = []
            st.rerun()

    st.divider()
    st.markdown("**현재 항목 목록** (직접 수정 / 노하우 입력 / 삭제)")
    categories_present = sorted(set(it["category"] for it in st.session_state.items))
    for cat in categories_present:
        with st.expander(f"■ {cat}", expanded=False):
            for it in [x for x in st.session_state.items if x["category"] == cat]:
                cols = st.columns([2, 2.5, 2.5, 0.7, 2, 0.5])
                it["item"] = cols[0].text_input("항목", it["item"], key=f"item_{it['id']}", label_visibility="collapsed")
                it["criteria"] = cols[1].text_input("판단기준", it["criteria"], key=f"crit_{it['id']}", label_visibility="collapsed")
                it["basis"] = cols[2].text_input("근거", it["basis"], key=f"basis_{it['id']}", label_visibility="collapsed")
                it["checklist"] = cols[3].checkbox("체크", it.get("checklist", True), key=f"chk_{it['id']}")
                it["user_note"] = cols[4].text_input("현장 노하우/비고", it.get("user_note", ""), key=f"note_{it['id']}", label_visibility="collapsed")
                if cols[5].button("✕", key=f"del_{it['id']}"):
                    st.session_state.items = [x for x in st.session_state.items if x["id"] != it["id"]]
                    st.rerun()

    with st.form("manual_add", clear_on_submit=True):
        st.markdown("**항목 직접 추가 (사용자 노하우 기반)**")
        c1, c2, c3, c4 = st.columns(4)
        n_cat = c1.text_input("카테고리", value="기타")
        n_item = c2.text_input("항목명")
        n_crit = c3.text_input("판단기준")
        n_basis = c4.text_input("근거")
        if st.form_submit_button("추가") and n_item:
            st.session_state.items.append({"id": new_id(), "category": n_cat, "item": n_item,
                                             "criteria": n_crit, "basis": n_basis, "checklist": True,
                                             "user_note": "", "focus": False})
            st.rerun()

    st.divider()
    c1, c2 = st.columns(2)
    if c1.button("← 이전"):
        goto(1)
    if c2.button("다음 단계 →", type="primary"):
        if not st.session_state.items:
            st.warning("최소 1개 이상의 항목이 필요합니다.")
        else:
            goto(3)

# ============================================================
# STAGE 3. 참고 서식 업로드 → 스타일 매칭
# ============================================================
elif st.session_state.stage == 3:
    st.subheader("STEP 3. 과거 표준서 참고 → 형식/양식 매칭")
    st.caption("과거에 사용하던 표준서(.docx)를 업로드하면 헤딩 폰트/색상/표 스타일을 분석해 새 문서에 동일하게 적용합니다. "
               "업로드하지 않으면 기본 사내 스타일(네이비 헤더)로 생성됩니다.")

    uploaded = st.file_uploader("참고 표준서 업로드 (.docx, 여러 개 가능)", type=["docx"], accept_multiple_files=True)

    if uploaded:
        profile = {}
        for f in uploaded:
            p = analyze_reference_docx(f.read())
            for k, v in p.items():
                if v and not profile.get(k):
                    profile[k] = v
        st.session_state.style_profile = profile
        st.success("참고 파일 스타일 분석 완료")
    style = merged_style(st.session_state.style_profile)

    st.markdown("**적용될 스타일 미리보기**")
    c1, c2, c3 = st.columns(3)
    c1.metric("헤딩 폰트", style["heading_font"])
    c2.metric("헤딩 색상", f"#{style['heading_color']}")
    c3.metric("표 스타일", style["table_style"])
    st.session_state.final_style = style

    st.divider()
    c1, c2 = st.columns(2)
    if c1.button("← 이전"):
        goto(2)
    if c2.button("다음 단계 →", type="primary"):
        goto(4)

# ============================================================
# STAGE 4. 집중포인트 지정 + 최종본 확정
# ============================================================
elif st.session_state.stage == 4:
    st.subheader("STEP 4. 집중관리 포인트 지정 → 최종본 확정")
    st.caption("v1 문서에서 특별히 강조하고 싶은 항목(현장 경험상 자주 이슈 나는 항목 등)에 체크하세요. "
               "체크된 항목은 워드/PPT/엑셀에서 별도로 강조 표시됩니다.")

    categories_present = sorted(set(it["category"] for it in st.session_state.items))
    for cat in categories_present:
        st.markdown(f"**■ {cat}**")
        for it in [x for x in st.session_state.items if x["category"] == cat]:
            c1, c2 = st.columns([0.06, 0.94])
            it["focus"] = c1.checkbox("", value=it.get("focus", False), key=f"focus_{it['id']}")
            c2.markdown(f"{'[집중] ' if it['focus'] else ''}**{it['item']}** — {it['criteria']}")

    st.divider()
    if st.button("AI 최종 검토 코멘트 받기"):
        with st.spinner("전체 항목 검토 중..."):
            st.session_state.final_review_text = ai_final_review(st.session_state.context, st.session_state.items)
    if st.session_state.final_review_text:
        st.info(st.session_state.final_review_text)

    st.divider()
    c1, c2 = st.columns(2)
    if c1.button("← 이전"):
        goto(3)
    if c2.button("최종본 확정 → 다운로드", type="primary"):
        goto(5)

# ============================================================
# STAGE 5. 다운로드
# ============================================================
elif st.session_state.stage == 5:
    st.subheader("STEP 5. 최종 문서 다운로드")
    st.success(f"'{st.session_state.category_name}' 표준서 — 총 {len(st.session_state.items)}개 항목 "
               f"(집중포인트 {sum(1 for it in st.session_state.items if it.get('focus'))}개)")

    style = st.session_state.get("final_style", DEFAULT_STYLE)
    context = st.session_state.context

    c1, c2, c3 = st.columns(3)
    with c1:
        if st.button("워드(.docx) 생성", use_container_width=True):
            data = build_docx(context, st.session_state.items, style)
            st.download_button("워드 다운로드", data=data,
                                file_name=f"{st.session_state.category_name}_{context.get('product','') or '제품'}.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                use_container_width=True)
    with c2:
        if st.button("PPT(.pptx) 생성", use_container_width=True):
            data = build_pptx(context, st.session_state.items, style)
            st.download_button("PPT 다운로드", data=data,
                                file_name=f"{st.session_state.category_name}_{context.get('product','') or '제품'}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True)
    with c3:
        if st.button("체크리스트(.xlsx) 생성", use_container_width=True):
            data = build_xlsx(context, st.session_state.items)
            st.download_button("엑셀 다운로드", data=data,
                                file_name=f"{st.session_state.category_name}_체크리스트.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                use_container_width=True)

    st.divider()
    if st.button("← 항목 다시 검토"):
        goto(4)
    if st.button("새 표준서 처음부터 시작"):
        for k in ["stage", "category_name", "category_purpose", "items", "next_id",
                  "ai_add_suggestions", "ai_remove_suggestions", "style_profile", "context", "final_review_text"]:
            st.session_state.pop(k, None)
        init_state()
        st.rerun()
